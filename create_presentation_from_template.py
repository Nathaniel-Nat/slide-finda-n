#!/usr/bin/env python3
"""
Generate presentation for Gráfica Finda N using the Slidesgo template.
Uses pptx library to edit existing template slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import shutil

# Copy template to output
shutil.copy(
    "templates/Defense of University Project by Slidesgo.pptx",
    "Apresentacao_TFC_Grafica_Finda_N.pptx"
)

prs = Presentation("Apresentacao_TFC_Grafica_Finda_N.pptx")

# Slide mapping: We'll use specific slides from the template for our content
# Template has 55 slides, we need 23

def clear_slide_content(slide):
    """Remove all content from a slide except background."""
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame:
            for paragraph in shape.text_frame.paragraphs:
                paragraph.clear()

def set_title(slide, title_text):
    """Set the title of a slide."""
    if slide.shapes.title:
        title_shape = slide.shapes.title
        title_shape.text = title_text
        # Format title
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(36)
                run.font.name = "Arial Black"

def set_content(slide, content_items):
    """Set bullet point content on a slide."""
    # Find the main content placeholder
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape != slide.shapes.title:
            tf = shape.text_frame
            tf.clear()
            
            for i, item in enumerate(content_items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = item
                p.level = 0
                p.font.size = Pt(16)
                p.font.name = "Calibri"
                
                # Add bullet
                p.bullet = True
            
            return

def set_two_column_content(slide, left_items, right_items):
    """Set content in two columns if available."""
    shapes = list(slide.shapes)
    
    # Try to find two text boxes for columns
    text_boxes = [s for s in shapes if hasattr(s, "text_frame")]
    
    if len(text_boxes) >= 2:
        # Left column
        tf_left = text_boxes[0].text_frame
        tf_left.clear()
        for i, item in enumerate(left_items):
            if i == 0:
                p = tf_left.paragraphs[0]
            else:
                p = tf_left.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(14)
        
        # Right column
        tf_right = text_boxes[1].text_frame
        tf_right.clear()
        for i, item in enumerate(right_items):
            if i == 0:
                p = tf_right.paragraphs[0]
            else:
                p = tf_right.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(14)

# Content for each slide based on the original outline
slides_content = [
    {
        "title": "Sistema de Gestão para a Gráfica Finda N",
        "subtitle": ["Trabalho de Fim de Curso", "Curso Médio Técnico de Informática", 
                     "Instituto Médio Técnico Luíza Andaluz", 
                     "Orientador: Prof. Adriano Jacinto"]
    },
    {
        "title": "Conhecendo a Gráfica Finda N",
        "content": [
            "Localizada no Rocha Pinto, Luanda",
            "Inaugurada em Março de 2025",
            "Especializada em materiais impressos personalizados",
            "Proprietário: Dr. Finda Henriques João Neto",
            "Equipa de 3 funcionários"
        ]
    },
    {
        "title": "O Problema",
        "content": [
            "Processos manuais: pedidos em papel, estoque de memória",
            "Informações frequentemente perdidas",
            "Atrasos no atendimento ao cliente",
            "Erros repetitivos e retrabalho constante",
            "Pergunta: Como um sistema pode melhorar os serviços?"
        ]
    },
    {
        "title": "Objetivos do Projeto",
        "content": [
            "Geral: Desenvolver sistema de gestão automatizado",
            "Analisar processos atuais e identificar melhorias",
            "Projetar arquitetura e interfaces do sistema",
            "Desenvolver funcionalidades principais",
            "Validar sistema com colaboradores da gráfica"
        ]
    },
    {
        "title": "Hipóteses do Projeto",
        "content": [
            "Automatizar tarefas manuais (pedidos, estoque, relatórios)",
            "Melhorar fluxo de trabalho entre setores",
            "Aumentar produtividade e atendimento",
            "Facilitar controle de estoque em tempo real",
            "Reduzir erros operacionais"
        ]
    },
    {
        "title": "O Que É um Sistema de Gestão?",
        "content": [
            "Recebe dados → processa → gera informação útil",
            "Na gráfica: entrada = pedidos/materiais; saída = produtos",
            "Permite planear, organizar, dirigir e controlar",
            "Feedback = avaliação do cliente",
            "Gestão eficiente de recursos"
        ]
    },
    {
        "title": "Tecnologias Web Utilizadas",
        "table": [
            ["Tecnologia", "Função"],
            ["HTML5", "Estrutura das páginas"],
            ["CSS3", "Visual e design responsivo"],
            ["JavaScript", "Interatividade e dinamismo"],
            ["PHP", "Lógica do servidor"],
            ["MySQL", "Armazenamento dos dados"]
        ]
    },
    {
        "title": "A Gráfica Como Organização",
        "content": [
            "Missão: Soluções gráficas inovadoras e de alta qualidade",
            "Visão: Ser líder no mercado gráfico regional",
            "Valores: Foco no cliente, qualidade e respeito pelo prazo",
            "Compromisso com excelência",
            "Atendimento personalizado"
        ]
    },
    {
        "title": "Metodologia de Trabalho",
        "content": [
            "Pesquisa aplicada com abordagem qualitativa",
            "Estudo de caso direto na Gráfica Finda N",
            "Amostra: 3 participantes (proprietário + 2 funcionários)",
            "Análise detalhada dos processos",
            "Validação contínua com utilizadores"
        ]
    },
    {
        "title": "Recolha de Informações",
        "content": [
            "Entrevistas semiestruturadas com proprietário e funcionários",
            "Temas: fluxo atual, dificuldades, expectativas",
            "Observação direta: 2 visitas de ~2 horas cada",
            "Acompanhamento do processo completo",
            "Identificação de gargalos e erros"
        ]
    },
    {
        "title": "Fases do Desenvolvimento",
        "content": [
            "1. Levantamento de requisitos (entrevistas + observação)",
            "2. Análise e definição do sistema",
            "3. Prototipagem no PowerPoint",
            "4. Desenvolvimento em HTML, CSS, JavaScript e PHP",
            "5. Testes e validação com utilizadores reais"
        ]
    },
    {
        "title": "Dificuldades Encontradas",
        "content": [
            "Informações desorganizadas na gráfica",
            "Ausência de registo digital anterior",
            "Tempo limitado para pesquisa",
            "Resistência inicial em partilhar dados",
            "Necessidade de adaptação contínua"
        ]
    },
    {
        "title": "Ferramentas Utilizadas",
        "content": [
            "Editor: Visual Studio Code",
            "Servidor local: XAMPP (Apache + MySQL + PHP)",
            "Prototipagem: Microsoft PowerPoint",
            "Razão: gratuitas, sem instalação complexa",
            "Adequadas à realidade angolana"
        ]
    },
    {
        "title": "Arquitetura do Sistema",
        "content": [
            "Arquitetura Cliente-Servidor",
            "Utilizador interage via navegador",
            "Servidor processa pedidos e comunica com BD",
            "Base de dados com 7 tabelas relacionais",
            "Tabela central: movimentos (atividade financeira)"
        ]
    },
    {
        "title": "Estrutura da Base de Dados",
        "content": [
            "Produtos → Stock → Encomendas",
            "Movimentos centraliza tudo: vendas, sessões, serviços",
            "Jogos e Serviços alimentam movimentos financeiros",
            "Normalizada até 3ª Forma Normal",
            "Relacionamentos bem definidos"
        ]
    },
    {
        "title": "Visão Geral do Sistema",
        "content": [
            "Acesso por qualquer navegador, sem instalação",
            "Interface em português, simples e intuitiva",
            "Tempo de aprendizagem: menos de 2 horas",
            "Módulos: Dashboard, Produtos, Stock, Jogos",
            "Serviços, Encomendas, Movimentos, Relatórios"
        ]
    },
    {
        "title": "Dashboard - Visão em Tempo Real",
        "content": [
            "Receita do dia atualizada automaticamente",
            "Máquinas ocupadas vs. disponíveis",
            "Encomendas pendentes com status",
            "Produtos com stock crítico alertados",
            "Indicadores chave de desempenho"
        ]
    },
    {
        "title": "Funcionalidades Principais",
        "content": [
            "Gestão de Produtos & Stock: alertas automáticos",
            "Gestão de Jogos: temporizador em tempo real",
            "Gestão de Encomendas: fluxo completo de status",
            "Serviços rápidos: registo em menos de 10 segundos",
            "Movimentos Financeiros: livro de caixa digital"
        ]
    },
    {
        "title": "Resultados Concretos",
        "table": [
            ["Métrica", "Resultado"],
            ["Tempo de venda de produto", "12 segundos / 3 cliques"],
            ["Tempo de início de sessão", "8 segundos"],
            ["Erros de cobrança", "0% (cálculo automático)"],
            ["Redução de tempo por transação", "~40%"],
            ["Tempo de carregamento", "< 2 segundos"]
        ]
    },
    {
        "title": "O Que o Sistema Prova",
        "content": [
            "Informatização viável para PMEs em Angola",
            "Tecnologias simples e gratuitas são suficientes",
            "Planeamento cuidadoso é fundamental",
            "Foco nas necessidades reais do utilizador",
            "Soluções de qualidade com recursos limitados"
        ]
    },
    {
        "title": "Limitações e Melhorias Futuras",
        "content": [
            "Senhas sem encriptação → implementar bcrypt",
            "Single-user → múltiplos operadores com permissões",
            "Relatórios básicos → exportação PDF/Excel",
            "Sem backup automático → integrar agendamento",
            "Gráficos de tendências por adicionar"
        ]
    },
    {
        "title": "Conclusão",
        "quote": "A implementação de um sistema de gestão informatizado representa uma solução eficaz para modernizar a Gráfica Finda N, reduzir erros, agilizar processos e elevar a qualidade do atendimento."
    },
    {
        "title": "Agradecimentos",
        "content": [
            "Orientador Prof. Adriano Jacinto",
            "Professores do IMTLA",
            "Dr. Finda Henriques João Neto e equipa",
            "Familiares e colegas",
            "Todos que contribuíram para este projeto"
        ]
    }
]

# Process slides
# Slide 1 is title slide - handle differently
slide1 = prs.slides[0]
if slide1.shapes.title:
    slide1.shapes.title.text = slides_content[0]["title"]
    # Format title
    for paragraph in slide1.shapes.title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.bold = True
            run.font.size = Pt(40)

# Add subtitle to slide 1
for shape in slide1.shapes:
    if hasattr(shape, "text_frame") and shape != slide1.shapes.title:
        tf = shape.text_frame
        tf.clear()
        for i, line in enumerate(slides_content[0]["subtitle"]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.alignment = 1  # Center
            p.font.size = Pt(18)

# Process remaining slides (2-23)
for i, content in enumerate(slides_content[1:], start=1):
    if i >= len(prs.slides):
        # Add new slide if needed
        layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(layout)
    else:
        slide = prs.slides[i]
    
    # Set title
    if "title" in content:
        if slide.shapes.title:
            slide.shapes.title.text = content["title"]
            for paragraph in slide.shapes.title.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.size = Pt(32)
    
    # Handle different content types
    if "content" in content:
        # Find content placeholder
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape != slide.shapes.title:
                tf = shape.text_frame
                tf.clear()
                for j, item in enumerate(content["content"]):
                    if j == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = item
                    p.level = 0
                    p.font.size = Pt(16)
                    p.font.name = "Calibri"
                    p.bullet = True
    
    elif "table" in content:
        # Create table
        table_data = content["table"]
        rows = len(table_data)
        cols = len(table_data[0]) if rows > 0 else 0
        
        # Find a good position for the table
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9.0)
        height = Inches(0.8) * rows
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Set column widths
        table.columns[0].width = Inches(3.5)
        table.columns[1].width = Inches(5.5)
        
        # Fill table
        for row_idx, row in enumerate(table_data):
            for col_idx, cell_text in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = cell_text
                
                # Format header row
                if row_idx == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0x0D, 0x94, 0x88)
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.bold = True
                            run.font.size = Pt(14)
                            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                else:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(12)
    
    elif "quote" in content:
        # Find content area and add quote
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape != slide.shapes.title:
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = content["quote"]
                p.alignment = 1  # Center
                p.font.size = Pt(18)
                p.font.italic = True

# Save presentation
prs.save("Apresentacao_TFC_Grafica_Finda_N.pptx")
print(f"Presentation created successfully with {len(prs.slides)} slides!")
print("Output file: Apresentacao_TFC_Grafica_Finda_N.pptx")
