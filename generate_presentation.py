#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Script para gerar os 23 slides da apresentação do TFC
Tema: Sistema de Gestão para a Gráfica Finda N
Template: Defense of University Project by Slidesgo
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Criar nova apresentação (do zero)
prs = Presentation()
output_path = "/workspace/Apresentacao_TFC_Grafica_Finda_N.pptx"

# Configurar tamanho do slide (16:9)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Definir estilo de texto consistente
def set_text_style(text_frame, font_size=14, bold=False, color=None):
    """Aplica estilo consistente ao texto"""
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            if color:
                run.font.color.rgb = color

def add_bullet_slide(prs, title, bullets, subtitle=None):
    """Adiciona um slide com título e lista de pontos"""
    # Usar layout de título e conteúdo
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # Definir título
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Estilizar título
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(24)
            run.font.bold = True
    
    # Adicionar subtítulo se existir
    if subtitle:
        left = Inches(0.5)
        top = Inches(1.2)
        width = Inches(9)
        height = Inches(0.5)
        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        subtitle_tf = subtitle_box.text_frame
        subtitle_tf.text = subtitle
        for paragraph in subtitle_tf.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(16)
                run.font.italic = True
    
    # Adicionar conteúdo (bullets)
    content_placeholder = slide.placeholders[1]
    tf = content_placeholder.text_frame
    tf.clear()  # Limpar conteúdo padrão
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        
        # Garantir tamanho mínimo de 14pt para todo o texto
        for run in p.runs:
            run.font.size = Pt(14)
        
        # Adicionar espaçamento entre parágrafos
        if i > 0:
            p.space_before = Pt(6)
    
    return slide

def add_title_slide(prs, title, subtitle, presenter=None):
    """Adiciona um slide de título"""
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Definir título
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Estilizar título
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(32)
            run.font.bold = True
    
    # Definir subtítulo
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle
    
    # Estilizar subtítulo
    for paragraph in subtitle_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(18)
    
    return slide

def add_table_slide(prs, title, headers, data, caption=None):
    """Adiciona um slide com tabela"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Título
    title_shape = slide.shapes.title
    title_shape.text = title
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(24)
            run.font.bold = True
    
    # Criar tabela
    rows = len(data) + 1
    cols = len(headers)
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(0.8)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Definir largura das colunas
    for i in range(cols):
        table.columns[i].width = Inches(8 / cols)
    
    # Preencher cabeçalhos
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(14)
                run.font.bold = True
    
    # Preencher dados
    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_data)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(13)
    
    # Adicionar legenda se existir
    if caption:
        left = Inches(1)
        top = Inches(1.5 + 0.8 + 0.3)
        width = Inches(8)
        height = Inches(0.5)
        caption_box = slide.shapes.add_textbox(left, top, width, height)
        caption_tf = caption_box.text_frame
        caption_tf.text = caption
        for paragraph in caption_tf.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(12)
                run.font.italic = True
    
    return slide

def add_quote_slide(prs, title, quote, author=None):
    """Adiciona um slide com citação"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Título
    title_shape = slide.shapes.title
    title_shape.text = title
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(24)
            run.font.bold = True
    
    # Adicionar citação
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(3)
    quote_box = slide.shapes.add_textbox(left, top, width, height)
    quote_tf = quote_box.text_frame
    quote_tf.word_wrap = True
    quote_tf.text = f'"{quote}"'
    
    for paragraph in quote_tf.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(18)
            run.font.italic = True
    
    # Adicionar autor se existir
    if author:
        author_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.5))
        author_tf = author_box.text_frame
        author_tf.text = f"— {author}"
        for paragraph in author_tf.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(14)
    
    return slide

# ============================================================================
# CONTEÚDO DOS SLIDES
# ============================================================================

print("A gerar slides...")

# Slide 1 — Capa
slide1 = add_title_slide(
    prs,
    "Sistema de Gestão para a Gráfica Finda N",
    "Trabalho de Fim de Curso\nCurso Médio Técnico de Informática\nInstituto Médio Técnico Luíza Andaluz\nOrientador: Prof. Adriano Jacinto"
)

# Slide 2 — Quem Somos e o Que Apresentamos
slide2 = add_bullet_slide(
    prs,
    "Quem Somos e o Que Apresentamos",
    [
        "Trabalho de Fim de Curso — Curso Médio Técnico de Informática",
        "Instituto Médio Técnico Luíza Andaluz",
        "Tema: Sistema de Gestão para a Gráfica Finda N",
        "Orientador: Prof. Adriano Jacinto",
        "Equipa: 6 apresentadores"
    ]
)

# Slide 3 — Conhecendo a Gráfica Finda N
slide3 = add_bullet_slide(
    prs,
    "Conhecendo a Gráfica Finda N",
    [
        "Localização: Rocha Pinto, Luanda",
        "Inauguração: Março de 2025",
        "Especialidade: Materiais impressos personalizados",
        "Proprietário: Dr. Finda Henriques João Neto",
        "Equipa: 3 funcionários"
    ]
)

# Slide 4 — O Problema
slide4 = add_bullet_slide(
    prs,
    "O Problema",
    [
        "Processos manuais: pedidos em papel, estoque de memória",
        "Informações frequentemente perdidas",
        "Atrasos no atendimento ao cliente",
        "Erros repetitivos e retrabalho constante",
        "Pergunta: Como um sistema pode melhorar os serviços?"
    ]
)

# Slide 5 — Objetivos do Projeto
slide5 = add_bullet_slide(
    prs,
    "Objetivos do Projeto",
    [
        "Geral: Desenvolver sistema de gestão para automatizar processos",
        "Analisar processos atuais e identificar melhorias",
        "Projetar arquitetura e interfaces do sistema",
        "Desenvolver funcionalidades principais",
        "Validar sistema com colaboradores da gráfica"
    ]
)

# Slide 6 — Hipóteses
slide6 = add_bullet_slide(
    prs,
    "O Que Esperávamos (Hipóteses)",
    [
        "Automatizar tarefas manuais (pedidos, estoque, relatórios)",
        "Melhorar fluxo de trabalho entre setores",
        "Aumentar produtividade e atendimento",
        "Facilitar controle de estoque em tempo real",
        "Reduzir erros operacionais"
    ]
)

# Slide 7 — O Que É um Sistema de Gestão?
slide7 = add_bullet_slide(
    prs,
    "O Que É um Sistema de Gestão?",
    [
        "Recebe dados → processa → gera informação útil",
        "Entrada: pedidos e materiais | Saída: produtos entregues",
        "Feedback: avaliação do cliente",
        "Permite planear, organizar, dirigir e controlar",
        "Aumenta eficiência operacional"
    ]
)

# Slide 8 — Por Que Tecnologias Web?
slide8 = add_table_slide(
    prs,
    "Por Que Tecnologias Web?",
    ["Tecnologia", "Função"],
    [
        ["HTML5", "Estrutura das páginas"],
        ["CSS3", "Visual e design responsivo"],
        ["JavaScript", "Interatividade e dinamismo"],
        ["PHP", "Lógica do servidor"],
        ["MySQL", "Armazenamento dos dados"]
    ],
    caption="Acessível por qualquer computador com navegador — sem instalação complexa"
)

# Slide 9 — A Gráfica Como Organização
slide9 = add_bullet_slide(
    prs,
    "A Gráfica Como Organização",
    [
        "Missão: Soluções gráficas inovadoras e de alta qualidade",
        "Visão: Ser líder no mercado gráfico regional",
        "Valores: Foco no cliente, qualidade e respeito pelo prazo",
        "Compromisso com satisfação do cliente",
        "Excelência em serviços impressos"
    ]
)

# Slide 10 — Como Trabalhámos
slide10 = add_bullet_slide(
    prs,
    "Como Trabalhámos (Metodologia)",
    [
        "Pesquisa aplicada com abordagem qualitativa",
        "Estudo de caso direto na Gráfica Finda N",
        "Amostra: 3 participantes (proprietário + 2 funcionários)",
        "Duração: 2 meses de pesquisa e desenvolvimento",
        "Foco nas necessidades reais dos utilizadores"
    ]
)

# Slide 11 — Recolha de Informações
slide11 = add_bullet_slide(
    prs,
    "Como Recolhemos as Informações",
    [
        "Entrevistas semiestruturadas com proprietário e funcionários",
        "Temas: fluxo atual, dificuldades, expectativas",
        "Observação direta: 2 visitas de ~2 horas cada",
        "Acompanhamento do processo completo",
        "Identificação de gargalos e erros operacionais"
    ]
)

# Slide 12 — Fases do Desenvolvimento
slide12 = add_bullet_slide(
    prs,
    "As 5 Fases do Desenvolvimento",
    [
        "1. Levantamento de requisitos (entrevistas + observação)",
        "2. Análise e definição de funcionalidades",
        "3. Prototipagem no PowerPoint para validação",
        "4. Desenvolvimento em HTML, CSS, JavaScript e PHP",
        "5. Testes e validação com utilizadores reais"
    ]
)

# Slide 13 — Dificuldades Encontradas
slide13 = add_bullet_slide(
    prs,
    "Dificuldades Encontradas",
    [
        "Informações desorganizadas na gráfica",
        "Ausência de registo digital anterior",
        "Tempo limitado para pesquisa",
        "Resistência inicial em partilhar dados",
        "Necessidade de adaptação contínua"
    ]
)

# Slide 14 — Ferramentas Utilizadas
slide14 = add_bullet_slide(
    prs,
    "Ferramentas Utilizadas",
    [
        "Editor: Visual Studio Code",
        "Servidor local: XAMPP (Apache + MySQL + PHP)",
        "Prototipagem: Microsoft PowerPoint",
        "Todas gratuitas e de fácil instalação",
        "Adequadas à realidade angolana"
    ]
)

# Slide 15 — Arquitetura do Sistema
slide15 = add_bullet_slide(
    prs,
    "Arquitetura do Sistema",
    [
        "Modelo Cliente-Servidor",
        "Utilizador interage via navegador",
        "Servidor processa pedidos e comunica com BD",
        "Base de dados: 7 tabelas relacionais (3ª Forma Normal)",
        "Tabela central: movimentos (atividade financeira)"
    ],
    subtitle="Ver Modelo Entidade-Relacionamento no próximo slide"
)

# Slide 16 — Estrutura da Base de Dados
slide16 = add_bullet_slide(
    prs,
    "Estrutura da Base de Dados",
    [
        "Produtos → ligados ao Stock → ligados às Encomendas",
        "Movimentos centraliza: vendas, sessões, serviços",
        "Jogos e Serviços alimentam movimentos financeiros",
        "Normalização até 3ª Forma Normal",
        "Relacionamentos bem definidos"
    ],
    subtitle="(Incluir imagem do Modelo Entidade-Relacionamento)"
)

# Slide 17 — Visão Geral do Sistema
slide17 = add_bullet_slide(
    prs,
    "Visão Geral do Sistema",
    [
        "Acesso por qualquer navegador, sem instalação",
        "Interface em português, simples e intuitiva",
        "Tempo de aprendizagem: < 2 horas",
        "Módulos: Dashboard, Produtos, Stock, Jogos",
        "Módulos: Serviços, Encomendas, Movimentos, Relatórios"
    ]
)

# Slide 18 — Dashboard
slide18 = add_bullet_slide(
    prs,
    "Dashboard — O Coração do Sistema",
    [
        "Visão em tempo real do negócio",
        "Receita do dia",
        "Máquinas ocupadas vs. disponíveis",
        "Encomendas pendentes",
        "Produtos com stock crítico"
    ],
    subtitle="(Incluir screenshot do Dashboard)"
)

# Slide 19 — Funcionalidades Principais
slide19 = add_bullet_slide(
    prs,
    "Funcionalidades Principais",
    [
        "Gestão de Produtos & Stock: alertas automáticos",
        "Gestão de Jogos: temporizador em tempo real",
        "Gestão de Encomendas: fluxo completo de status",
        "Serviços rápidos: registo em < 10 segundos",
        "Movimentos Financeiros: livro de caixa digital"
    ]
)

# Slide 20 — Resultados Concretos
slide20 = add_table_slide(
    prs,
    "Resultados Concretos",
    ["Métrica", "Resultado"],
    [
        ["Tempo de venda de produto", "12 segundos / 3 cliques"],
        ["Tempo de início de sessão", "8 segundos"],
        ["Erros de cobrança", "0% (cálculo automático)"],
        ["Redução de tempo por transação", "~40%"],
        ["Tempo de carregamento", "< 2 segundos"]
    ]
)

# Slide 21 — Conclusões
slide21 = add_bullet_slide(
    prs,
    "O Que o Sistema Prova",
    [
        "Informatização viável para PMEs em Angola",
        "Tecnologias simples criam soluções de qualidade",
        "Planeamento cuidadoso é fundamental",
        "Foco nas necessidades do utilizador > tecnologia",
        "Soluções locais podem ser eficazes"
    ]
)

# Slide 22 — Limitações e Futuro
slide22 = add_bullet_slide(
    prs,
    "Limitações e Caminho à Frente",
    [
        "Senhas sem encriptação → implementar bcrypt",
        "Um utilizador → múltiplos operadores com permissões",
        "Relatórios básicos → exportação PDF/Excel",
        "Sem backup automático → integrar agendamento",
        "Melhoria contínua garantida"
    ]
)

# Slide 23 — Conclusão Final
slide23 = add_quote_slide(
    prs,
    "Conclusão",
    "A implementação de um sistema de gestão informatizado representa uma solução eficaz para modernizar a Gráfica Finda N, reduzir erros, agilizar processos e elevar a qualidade do atendimento.",
    author="Trabalho de Fim de Curso - IMTLA"
)

# Guardar apresentação
prs.save(output_path)
print(f"Apresentação guardada em: {output_path}")
print(f"Total de slides: {len(prs.slides)}")
